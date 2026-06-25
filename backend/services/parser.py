from typing import List, Dict, Any
import os
import io
import html
import re
import fitz
from docx import Document
from bs4 import BeautifulSoup
from PIL import Image
from PIL.ExifTags import TAGS


def parse_pdf(file_path: str) -> List[Dict[str, Any]]:
    doc = fitz.open(file_path)
    chunks = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text()
        if text.strip():
            chunks.append({
                "content": text.strip(),
                "meta": {"page": page_num + 1, "type": "pdf"}
            })
        # 提取页面中的图片并做简单 OCR（如果可用）
        try:
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list, start=1):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                ext = base_image["ext"]
                img_result = _ocr_image_bytes(image_bytes, ext)
                if img_result:
                    chunks.append({
                        "content": f"[页面 {page_num + 1} 图片 {img_index}]\n{img_result}",
                        "meta": {"page": page_num + 1, "type": "pdf_image", "image_index": img_index}
                    })
        except Exception:
            pass
    return chunks


def parse_word(file_path: str) -> List[Dict[str, Any]]:
    doc = Document(file_path)
    chunks = []
    for i, para in enumerate(doc.paragraphs, 1):
        if para.text.strip():
            chunks.append({
                "content": para.text.strip(),
                "meta": {"paragraph": i, "type": "docx"}
            })
    # 提取表格
    for ti, table in enumerate(doc.tables, 1):
        rows = []
        for row in table.rows:
            rows.append(" | ".join(cell.text.strip() for cell in row.cells))
        if rows:
            chunks.append({
                "content": f"[表格 {ti}]\n" + "\n".join(rows),
                "meta": {"table": ti, "type": "docx_table"}
            })
    return chunks


def parse_txt(file_path: str) -> List[Dict[str, Any]]:
    with open(file_path, 'r', encoding='utf-8') as f:
        text = f.read()
    chunks = []
    paragraphs = text.split('\n\n')
    for i, p in enumerate(paragraphs, 1):
        if p.strip():
            chunks.append({
                "content": p.strip(),
                "meta": {"paragraph": i, "type": "txt"}
            })
    return chunks


def parse_markdown(file_path: str) -> List[Dict[str, Any]]:
    try:
        import markdown
    except ImportError:
        raise ValueError("未安装 markdown，无法解析 Markdown")

    with open(file_path, 'r', encoding='utf-8') as f:
        md_text = f.read()
    # 简单分段：按标题和空行拆分
    import re
    sections = re.split(r'\n(?=#+\s)', md_text)
    chunks = []
    for i, section in enumerate(sections, 1):
        section = section.strip()
        if section:
            chunks.append({
                "content": section,
                "meta": {"section": i, "type": "markdown"}
            })
    return chunks


def parse_pptx(file_path: str) -> List[Dict[str, Any]]:
    try:
        from pptx import Presentation
    except ImportError:
        raise ValueError("未安装 python-pptx，无法解析 PPT")

    prs = Presentation(file_path)
    chunks = []
    for si, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            chunks.append({
                "content": f"[幻灯片 {si}]\n" + "\n".join(texts),
                "meta": {"slide": si, "type": "pptx"}
            })
        # 尝试提取幻灯片中的图片
        try:
            for img_idx, shape in enumerate(slide.shapes, 1):
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    image = shape.image
                    img_result = _ocr_image_bytes(image.blob, image.ext)
                    if img_result:
                        chunks.append({
                            "content": f"[幻灯片 {si} 图片 {img_idx}]\n{img_result}",
                            "meta": {"slide": si, "type": "pptx_image", "image_index": img_idx}
                        })
        except Exception:
            pass
    return chunks


def parse_excel(file_path: str) -> List[Dict[str, Any]]:
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == '.xls':
            import pandas as pd
            df_dict = pd.read_excel(file_path, sheet_name=None, engine='xlrd')
        else:
            import pandas as pd
            df_dict = pd.read_excel(file_path, sheet_name=None, engine='openpyxl')
    except ImportError as e:
        raise ValueError(f"缺少解析依赖: {e}")
    except Exception as e:
        raise ValueError(f"Excel 解析失败: {e}")

    chunks = []
    for sheet_name, df in df_dict.items():
        df = df.fillna('')
        rows = []
        for row in df.itertuples(index=False, name=None):
            row_text = ' | '.join(str(cell) if cell is not None and str(cell) != '' else '' for cell in row)
            row_text = row_text.rstrip(' |').strip()
            if row_text:
                rows.append(row_text)
        if rows:
            chunks.append({
                "content": f"[工作表 {sheet_name}]\n" + "\n".join(rows),
                "meta": {"sheet": sheet_name, "type": ext.lstrip('.')}
            })
    return chunks


def parse_csv(file_path: str) -> List[Dict[str, Any]]:
    import csv
    chunks = []
    with open(file_path, 'r', encoding='utf-8-sig', newline='') as f:
        reader = csv.reader(f)
        rows = [" | ".join(cell.strip() for cell in row) for row in reader if any(cell.strip() for cell in row)]
    if rows:
        chunks.append({
            "content": "\n".join(rows),
            "meta": {"type": "csv"}
        })
    return chunks


def _ocr_image_bytes(image_bytes: bytes, ext: str) -> str:
    """尝试 OCR；未安装 tesseract 或识别失败时返回空字符串。"""
    try:
        from pytesseract import image_to_string
        img = Image.open(io.BytesIO(image_bytes))
        # 转换为 RGB 避免模式问题
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        text = image_to_string(img, lang="chi_sim+eng")
        return text.strip()
    except Exception:
        return ""


def parse_image(file_path: str) -> List[Dict[str, Any]]:
    """解析图片：优先 OCR，失败则提取 EXIF/尺寸等元数据。"""
    try:
        img = Image.open(file_path)
    except Exception as e:
        raise ValueError(f"无法打开图片: {e}")

    result = []
    ocr_text = _ocr_image_path(file_path)
    if ocr_text:
        result.append({
            "content": f"[图片 OCR 识别结果]\n{ocr_text}",
            "meta": {"type": "image_ocr", "width": img.width, "height": img.height}
        })

    meta_parts = [f"尺寸: {img.width} x {img.height}", f"格式: {img.format or '未知'}"]
    try:
        exif = img._getexif()
        if exif:
            for tag_id, value in exif.items():
                tag = TAGS.get(tag_id, tag_id)
                if tag in ("DateTime", "Make", "Model"):
                    meta_parts.append(f"{tag}: {value}")
    except Exception:
        pass

    result.append({
        "content": "[图片元数据]\n" + "\n".join(meta_parts),
        "meta": {"type": "image_meta", "width": img.width, "height": img.height}
    })
    return result


def _ocr_image_path(file_path: str) -> str:
    try:
        from pytesseract import image_to_string
        img = Image.open(file_path)
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        return image_to_string(img, lang="chi_sim+eng").strip()
    except Exception:
        return ""


def parse_online_html(file_path: str) -> List[Dict[str, Any]]:
    with open(file_path, 'r', encoding='utf-8') as f:
        html_text = f.read()
    soup = BeautifulSoup(html_text, 'html.parser')
    text = soup.get_text('\n')
    # 按空行分段，过滤过短段落
    paragraphs = [p.strip() for p in re.split(r'\n\s*\n', text) if p.strip()]
    chunks = []
    for i, p in enumerate(paragraphs, 1):
        chunks.append({
            "content": p,
            "meta": {"paragraph": i, "type": "online"}
        })
    if not chunks:
        chunks.append({
            "content": text.strip() or "[空文档]",
            "meta": {"paragraph": 1, "type": "online"}
        })
    return chunks


async def parse_file(file_path: str, file_type: str) -> List[Dict[str, Any]]:
    ft = file_type.lower()
    if ft == 'pdf':
        return parse_pdf(file_path)
    elif ft in ['doc', 'docx']:
        return parse_word(file_path)
    elif ft in ['txt', 'md', 'markdown']:
        return parse_markdown(file_path) if ft in ['md', 'markdown'] else parse_txt(file_path)
    elif ft in ['ppt', 'pptx']:
        return parse_pptx(file_path)
    elif ft in ['xls', 'xlsx']:
        return parse_excel(file_path)
    elif ft == 'csv':
        return parse_csv(file_path)
    elif ft == 'online':
        return parse_online_html(file_path)
    elif ft in ['png', 'jpg', 'jpeg', 'bmp', 'gif', 'webp', 'tiff']:
        return parse_image(file_path)
    else:
        raise ValueError(f"不支持的文件类型: {file_type}")
