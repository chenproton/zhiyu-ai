import html
import os
from fastapi.responses import HTMLResponse, FileResponse
from fastapi import HTTPException
from docx import Document


def generate_preview(file_path: str, file_type: str):
    ft = file_type.lower()
    if ft == 'pdf':
        # 前端会直接用 /download 接口取 blob 预览，这里保留 PDF 原文件返回
        return FileResponse(file_path, media_type="application/pdf")
    elif ft in ['doc', 'docx']:
        return _preview_word(file_path)
    elif ft in ['txt', 'md', 'markdown']:
        return _preview_markdown(file_path) if ft in ['md', 'markdown'] else _preview_txt(file_path)
    elif ft == 'online':
        return _preview_online(file_path)
    elif ft in ['ppt', 'pptx']:
        return _preview_pptx(file_path)
    elif ft in ['xls', 'xlsx', 'csv']:
        return _preview_excel(file_path)
    elif ft in ['png', 'jpg', 'jpeg', 'bmp', 'gif', 'webp', 'tiff']:
        return _preview_image(file_path)
    else:
        raise HTTPException(400, "该文件类型暂不支持在线预览，请下载查看")


def _preview_word(file_path: str):
    try:
        doc = Document(file_path)
    except Exception as e:
        raise HTTPException(400, f"Word 文档解析失败: {e}")
    html_parts = ['<div style="font-family: Arial, sans-serif; line-height: 1.6; padding: 20px;">']
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            html_parts.append(f"<p>{html.escape(text)}</p>")
    # 表格
    for table in doc.tables:
        html_parts.append('<table style="border-collapse:collapse;width:100%;margin:12px 0">')
        for row in table.rows:
            html_parts.append('<tr>')
            for cell in row.cells:
                html_parts.append(f'<td style="border:1px solid #ddd;padding:8px">{html.escape(cell.text.strip())}</td>')
            html_parts.append('</tr>')
        html_parts.append('</table>')
    html_parts.append('</div>')
    return HTMLResponse(content="".join(html_parts))


def _preview_txt(file_path: str):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
    except Exception as e:
        raise HTTPException(400, f"文本读取失败: {e}")
    escaped = html.escape(content)
    return HTMLResponse(content=f'<pre style="white-space: pre-wrap; padding: 20px;">{escaped}</pre>')


def _preview_online(file_path: str):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
    except Exception as e:
        raise HTTPException(400, f"在线文档读取失败: {e}")
    return HTMLResponse(content=f'<div style="font-family: Arial, sans-serif; line-height: 1.6; padding: 20px;" class="online-doc-preview">{content}</div>')


def _preview_markdown(file_path: str):
    try:
        import markdown
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
    except ImportError:
        raise HTTPException(400, "未安装 markdown，无法预览 Markdown")
    except Exception as e:
        raise HTTPException(400, f"Markdown 读取失败: {e}")
    md = markdown.Markdown(extensions=['extra', 'codehilite', 'toc'])
    rendered = md.convert(content)
    return HTMLResponse(content=f'<div style="font-family: Arial, sans-serif; line-height: 1.6; padding: 20px;" class="markdown-body">{rendered}</div>')


def _preview_pptx(file_path: str):
    try:
        from pptx import Presentation
    except ImportError:
        raise HTTPException(400, "未安装 python-pptx，无法预览 PPT")

    try:
        prs = Presentation(file_path)
    except Exception as e:
        raise HTTPException(400, f"PPT 解析失败: {e}")

    html_parts = ['<div style="font-family: Arial, sans-serif; line-height: 1.6; padding: 20px;">']
    for si, slide in enumerate(prs.slides, 1):
        html_parts.append(f'<div style="margin-bottom:24px;padding:16px;border:1px solid #e2e8f0;border-radius:8px">')
        html_parts.append(f'<div style="font-weight:bold;color:#4f46e5;margin-bottom:8px">幻灯片 {si}</div>')
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                html_parts.append(f"<p>{html.escape(shape.text.strip())}</p>")
        html_parts.append('</div>')
    html_parts.append('</div>')
    return HTMLResponse(content="".join(html_parts))


def _preview_excel(file_path: str):
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == '.csv':
            import pandas as pd
            df_dict = {'Sheet1': pd.read_csv(file_path, encoding='utf-8-sig')}
        elif ext == '.xls':
            import pandas as pd
            df_dict = pd.read_excel(file_path, sheet_name=None, engine='xlrd')
        else:
            import pandas as pd
            df_dict = pd.read_excel(file_path, sheet_name=None, engine='openpyxl')
    except ImportError as e:
        raise HTTPException(400, f"缺少解析依赖: {e}")
    except Exception as e:
        raise HTTPException(400, f"Excel/CSV 解析失败: {e}")

    html_parts = ['<div style="font-family: Arial, sans-serif; line-height: 1.6; padding: 20px;">']
    for sheet_name, df in df_dict.items():
        df = df.fillna('')
        html_parts.append(f'<div style="margin-bottom:24px">')
        html_parts.append(f'<div style="font-weight:bold;color:#4f46e5;margin-bottom:8px">工作表：{html.escape(str(sheet_name))}</div>')
        html_parts.append('<table style="border-collapse:collapse;width:100%">')
        for row in df.itertuples(index=False, name=None):
            html_parts.append('<tr>')
            for cell in row:
                val = "" if cell is None else str(cell)
                html_parts.append(f'<td style="border:1px solid #ddd;padding:8px">{html.escape(val)}</td>')
            html_parts.append('</tr>')
        html_parts.append('</table>')
        html_parts.append('</div>')
    html_parts.append('</div>')
    return HTMLResponse(content="".join(html_parts))


def _preview_image(file_path: str):
    # 图片直接返回文件让浏览器显示
    return FileResponse(file_path)
